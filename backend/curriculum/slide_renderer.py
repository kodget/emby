"""
Slide Rendering System
Converts PDF/PPTX/DOCX pages to images with text for the reader view.
Results are cached in SlideContent to avoid re-processing on every request.
"""

import io
import requests
import fitz  # PyMuPDF
import cloudinary.uploader
from PIL import Image


def _upload_page_image(image_bytes, slide_id, page_num):
    try:
        result = cloudinary.uploader.upload(
            image_bytes,
            folder="emby/slide_pages",
            public_id=f"slides/{slide_id}/page_{page_num}",
            resource_type='image',
            format='jpg',
            quality='auto:best',
            timeout=60,
        )
        return result['secure_url']
    except Exception as e:
        print(f"Cloudinary upload failed (page {page_num}): {e}")
        return None


def _render_pdf_from_bytes(pdf_bytes, slide_id):
    """Render PDF pages as images using pdf2image for high quality."""
    from pdf2image import convert_from_bytes
    
    print(f"Converting PDF to images...")
    
    try:
        # Convert PDF to images at 150 DPI (good balance of quality and size)
        images = convert_from_bytes(
            pdf_bytes,
            dpi=150,
            fmt='jpeg',
            thread_count=2
        )
        
        print(f"Converted {len(images)} pages")
        
        pages = []
        for page_num, img in enumerate(images, 1):
            # Convert to JPEG bytes
            buf = io.BytesIO()
            img.save(buf, format='JPEG', quality=85, optimize=True)
            
            print(f"Uploading page {page_num}...")
            image_url = _upload_page_image(buf.getvalue(), slide_id, page_num)
            if not image_url:
                print(f"Upload failed for page {page_num}")
                continue
            
            pages.append({
                "page_number": page_num,
                "image_url": image_url,
                "width": img.width,
                "height": img.height,
                "text_blocks": [],
            })
        
        return pages
        
    except Exception as e:
        print(f"pdf2image conversion failed: {e}")
        print("Falling back to PyMuPDF...")
        
        # Fallback to PyMuPDF if pdf2image fails
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        pages = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            zoom = 2
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, alpha=False)

            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            buf = io.BytesIO()
            img.save(buf, format='JPEG', quality=88, optimize=True)

            image_url = _upload_page_image(buf.getvalue(), slide_id, page_num + 1)
            if not image_url:
                continue

            pages.append({
                "page_number": page_num + 1,
                "image_url": image_url,
                "width": pix.width,
                "height": pix.height,
                "text_blocks": [],
            })
        doc.close()
        return pages


def _render_pptx_from_bytes(pptx_bytes, slide_id):
    """
    Render PPTX slides by converting to PDF first using LibreOffice, then rendering as images.
    Falls back to basic text extraction if LibreOffice is not available.
    """
    import tempfile
    import os
    import subprocess
    from pptx import Presentation
    
    print(f"Attempting PPTX rendering for {slide_id}...")
    
    # Try LibreOffice conversion first (best quality)
    try:
        # Save PPTX to temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_pptx:
            tmp_pptx.write(pptx_bytes)
            tmp_pptx_path = tmp_pptx.name
        
        # Create temp directory for output
        tmp_dir = tempfile.mkdtemp()
        
        try:
            # Convert PPTX to PDF using LibreOffice
            print(f"Converting PPTX to PDF using LibreOffice...")
            result = subprocess.run([
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', tmp_dir,
                tmp_pptx_path
            ], capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0:
                # Find the generated PDF
                pdf_path = os.path.join(tmp_dir, os.path.splitext(os.path.basename(tmp_pptx_path))[0] + '.pdf')
                
                if os.path.exists(pdf_path):
                    print(f"PDF conversion successful, rendering pages...")
                    with open(pdf_path, 'rb') as pdf_file:
                        pdf_bytes = pdf_file.read()
                    
                    # Use PDF rendering function
                    pages = _render_pdf_from_bytes(pdf_bytes, slide_id)
                    
                    # Cleanup
                    os.remove(pdf_path)
                    os.remove(tmp_pptx_path)
                    os.rmdir(tmp_dir)
                    
                    print(f"PPTX rendered via LibreOffice: {len(pages)} pages")
                    return pages
            else:
                print(f"LibreOffice conversion failed: {result.stderr}")
        
        except (subprocess.TimeoutExpired, FileNotFoundError) as e:
            print(f"LibreOffice not available or timeout: {e}")
        
        finally:
            # Cleanup temp files
            if os.path.exists(tmp_pptx_path):
                os.remove(tmp_pptx_path)
            if os.path.exists(tmp_dir):
                for file in os.listdir(tmp_dir):
                    os.remove(os.path.join(tmp_dir, file))
                os.rmdir(tmp_dir)
    
    except Exception as e:
        print(f"LibreOffice conversion error: {e}")
    
    # Fallback: Extract text and create simple rendered slides with text overlay
    print("Falling back to text extraction with visual rendering...")
    try:
        from PIL import ImageDraw, ImageFont
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        prs = Presentation(io.BytesIO(pptx_bytes))
        print(f"PPTX has {len(prs.slides)} slides")
        
        pages = []
        target_width = 1440
        target_height = 1080
        
        # Try to load a font for text rendering
        try:
            font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36)
            font_body = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        except:
            # Fallback to default font
            font_title = ImageFont.load_default()
            font_body = ImageFont.load_default()
        
        for slide_num, slide in enumerate(prs.slides, 1):
            try:
                # Create white canvas with light gray background
                canvas = Image.new('RGB', (target_width, target_height), '#f5f5f5')
                draw = ImageDraw.Draw(canvas)
                
                # Add a subtle border
                draw.rectangle([10, 10, target_width-10, target_height-10], outline='#cccccc', width=2)
                
                # Extract and render text from shapes
                y_offset = 60
                text_blocks = []
                
                for shape in slide.shapes:
                    if not hasattr(shape, "text") or not shape.text.strip():
                        continue
                    
                    text = shape.text.strip()
                    
                    # Determine if this is likely a title (first text box or larger font)
                    is_title = len(text_blocks) == 0 or len(text) < 100
                    
                    # Choose font and color
                    current_font = font_title if is_title else font_body
                    text_color = '#1a1a1a' if is_title else '#333333'
                    
                    # Word wrap text
                    max_width = target_width - 120
                    words = text.split()
                    lines = []
                    current_line = []
                    
                    for word in words:
                        test_line = ' '.join(current_line + [word])
                        # Approximate width check
                        if len(test_line) * (18 if is_title else 12) < max_width:
                            current_line.append(word)
                        else:
                            if current_line:
                                lines.append(' '.join(current_line))
                            current_line = [word]
                    
                    if current_line:
                        lines.append(' '.join(current_line))
                    
                    # Draw each line
                    for line in lines:
                        if y_offset < target_height - 100:  # Leave space at bottom
                            draw.text((60, y_offset), line, fill=text_color, font=current_font)
                            y_offset += 50 if is_title else 35
                    
                    y_offset += 20  # Space between text blocks
                    
                    # Store text block info
                    text_blocks.append({
                        "text": text,
                        "x": 60,
                        "y": y_offset - (len(lines) * (50 if is_title else 35)),
                        "width": max_width,
                        "height": len(lines) * (50 if is_title else 35),
                        "font_size": 36 if is_title else 24,
                    })
                
                # Add page number at bottom
                draw.text((target_width - 100, target_height - 40), 
                         f"Slide {slide_num}", fill='#999999', font=font_body)
                
                # Convert to bytes
                buf = io.BytesIO()
                canvas.save(buf, format='JPEG', quality=90, optimize=True)
                
                print(f"Uploading slide {slide_num}...")
                image_url = _upload_page_image(buf.getvalue(), slide_id, slide_num)
                if not image_url:
                    print(f"Upload failed for slide {slide_num}")
                    continue
                
                print(f"Slide {slide_num} uploaded: {image_url}")
                pages.append({
                    "page_number": slide_num,
                    "image_url": image_url,
                    "width": target_width,
                    "height": target_height,
                    "text_blocks": text_blocks,
                })
                
            except Exception as e:
                print(f"Error rendering slide {slide_num}: {e}")
                import traceback
                traceback.print_exc()
                continue
        
        print(f"PPTX rendering complete: {len(pages)} pages")
        return pages
        
    except Exception as e:
        print(f"Failed to render PPTX: {e}")
        import traceback
        traceback.print_exc()
        return []


def _render_docx_from_bytes(docx_bytes, slide_id):
    """
    Render DOCX as pages with actual text rendered visually.
    Uses LibreOffice for best quality, falls back to text rendering.
    """
    import tempfile
    import os
    import subprocess
    from docx import Document
    from PIL import ImageDraw, ImageFont
    
    print(f"Attempting DOCX rendering for {slide_id}...")
    
    # Try LibreOffice conversion first (best quality)
    try:
        # Save DOCX to temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_docx:
            tmp_docx.write(docx_bytes)
            tmp_docx_path = tmp_docx.name
        
        # Create temp directory for output
        tmp_dir = tempfile.mkdtemp()
        
        try:
            # Convert DOCX to PDF using LibreOffice
            print(f"Converting DOCX to PDF using LibreOffice...")
            result = subprocess.run([
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', tmp_dir,
                tmp_docx_path
            ], capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0:
                # Find the generated PDF
                pdf_path = os.path.join(tmp_dir, os.path.splitext(os.path.basename(tmp_docx_path))[0] + '.pdf')
                
                if os.path.exists(pdf_path):
                    print(f"PDF conversion successful, rendering pages...")
                    with open(pdf_path, 'rb') as pdf_file:
                        pdf_bytes = pdf_file.read()
                    
                    # Use PDF rendering function
                    pages = _render_pdf_from_bytes(pdf_bytes, slide_id)
                    
                    # Cleanup
                    os.remove(pdf_path)
                    os.remove(tmp_docx_path)
                    os.rmdir(tmp_dir)
                    
                    print(f"DOCX rendered via LibreOffice: {len(pages)} pages")
                    return pages
            else:
                print(f"LibreOffice conversion failed: {result.stderr}")
        
        except (subprocess.TimeoutExpired, FileNotFoundError) as e:
            print(f"LibreOffice not available or timeout: {e}")
        
        finally:
            # Cleanup temp files
            if os.path.exists(tmp_docx_path):
                os.remove(tmp_docx_path)
            if os.path.exists(tmp_dir):
                for file in os.listdir(tmp_dir):
                    os.remove(os.path.join(tmp_dir, file))
                os.rmdir(tmp_dir)
    
    except Exception as e:
        print(f"LibreOffice conversion error: {e}")
    
    # Fallback: Extract text and render it visually
    print("Falling back to text extraction with visual rendering...")
    try:
        doc = Document(io.BytesIO(docx_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        if not paragraphs:
            print("No text content found in DOCX")
            return []
        
        # Try to load fonts
        try:
            font_heading = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
            font_body = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 20)
        except:
            font_heading = ImageFont.load_default()
            font_body = ImageFont.load_default()
        
        # Calculate pages needed
        line_height = 32
        lines_per_page = 30
        img_width = 1200
        img_height = 1600
        
        pages = []
        current_page_lines = []
        text_blocks = []
        
        for para in paragraphs:
            # Word wrap paragraph
            words = para.split()
            current_line = []
            
            for word in words:
                test_line = ' '.join(current_line + [word])
                # Approximate width (adjust based on font)
                if len(test_line) * 10 < img_width - 120:
                    current_line.append(word)
                else:
                    if current_line:
                        current_page_lines.append(' '.join(current_line))
                    current_line = [word]
            
            if current_line:
                current_page_lines.append(' '.join(current_line))
            
            # Add blank line after paragraph
            current_page_lines.append('')
            
            # Check if we need a new page
            if len(current_page_lines) >= lines_per_page:
                # Render current page
                page_num = len(pages) + 1
                canvas = Image.new('RGB', (img_width, img_height), 'white')
                draw = ImageDraw.Draw(canvas)
                
                # Add border
                draw.rectangle([20, 20, img_width-20, img_height-20], outline='#dddddd', width=1)
                
                # Draw text
                y_pos = 60
                for line in current_page_lines[:lines_per_page]:
                    if line.strip():
                        # Check if it looks like a heading (short, all caps, or starts with number)
                        is_heading = len(line) < 60 and (line.isupper() or line[0].isdigit())
                        current_font = font_heading if is_heading else font_body
                        draw.text((60, y_pos), line, fill='#1a1a1a', font=current_font)
                    y_pos += line_height
                
                # Add page number
                draw.text((img_width - 100, img_height - 40), 
                         f"Page {page_num}", fill='#999999', font=font_body)
                
                # Upload page
                buf = io.BytesIO()
                canvas.save(buf, format='JPEG', quality=90, optimize=True)
                
                image_url = _upload_page_image(buf.getvalue(), slide_id, page_num)
                if image_url:
                    pages.append({
                        "page_number": page_num,
                        "image_url": image_url,
                        "width": img_width,
                        "height": img_height,
                        "text_blocks": [{"text": '\n'.join(current_page_lines[:lines_per_page]), 
                                       "x": 60, "y": 60, "width": img_width-120, 
                                       "height": lines_per_page * line_height, "font_size": 20}],
                    })
                
                # Start new page with remaining lines
                current_page_lines = current_page_lines[lines_per_page:]
        
        # Render final page if there are remaining lines
        if current_page_lines:
            page_num = len(pages) + 1
            canvas = Image.new('RGB', (img_width, img_height), 'white')
            draw = ImageDraw.Draw(canvas)
            
            draw.rectangle([20, 20, img_width-20, img_height-20], outline='#dddddd', width=1)
            
            y_pos = 60
            for line in current_page_lines:
                if line.strip():
                    is_heading = len(line) < 60 and (line.isupper() or line[0].isdigit())
                    current_font = font_heading if is_heading else font_body
                    draw.text((60, y_pos), line, fill='#1a1a1a', font=current_font)
                y_pos += line_height
            
            draw.text((img_width - 100, img_height - 40), 
                     f"Page {page_num}", fill='#999999', font=font_body)
            
            buf = io.BytesIO()
            canvas.save(buf, format='JPEG', quality=90, optimize=True)
            
            image_url = _upload_page_image(buf.getvalue(), slide_id, page_num)
            if image_url:
                pages.append({
                    "page_number": page_num,
                    "image_url": image_url,
                    "width": img_width,
                    "height": img_height,
                    "text_blocks": [{"text": '\n'.join(current_page_lines), 
                                   "x": 60, "y": 60, "width": img_width-120, 
                                   "height": len(current_page_lines) * line_height, "font_size": 20}],
                })
        
        print(f"DOCX rendering complete: {len(pages)} pages")
        return pages
        
    except Exception as e:
        print(f"Failed to render DOCX: {e}")
        import traceback
        traceback.print_exc()
        return []


def render_slide_pages(file_url, file_type, slide_id):
    """
    Download the file and render all pages.
    Returns {"total_pages": N, "pages": [...]} for the SlideContent cache.
    """
    print(f"Rendering {slide_id} ({file_type}) from {file_url}")
    try:
        resp = requests.get(file_url, timeout=30)
        if resp.status_code != 200 or not resp.content:
            print(f"Download failed: HTTP {resp.status_code}")
            return {"total_pages": 0, "pages": []}
        file_bytes = resp.content
        print(f"Downloaded {len(file_bytes)} bytes")
    except Exception as e:
        print(f"Download error: {e}")
        return {"total_pages": 0, "pages": []}

    ft = file_type.lower()
    try:
        if 'pdf' in ft:
            print("Rendering as PDF...")
            pages = _render_pdf_from_bytes(file_bytes, slide_id)
        elif 'pptx' in ft or 'presentation' in ft or 'ppt' in ft:
            print("Rendering as PPTX...")
            pages = _render_pptx_from_bytes(file_bytes, slide_id)
        elif 'docx' in ft or 'document' in ft or 'doc' in ft:
            print("Rendering as DOCX...")
            pages = _render_docx_from_bytes(file_bytes, slide_id)
        else:
            print(f"Unsupported file type: {file_type}")
            pages = []
        
        print(f"Rendered {len(pages)} pages")
    except Exception as e:
        import traceback
        print(f"Rendering error for {slide_id}: {e}")
        traceback.print_exc()
        pages = []

    return {"total_pages": len(pages), "pages": pages}
