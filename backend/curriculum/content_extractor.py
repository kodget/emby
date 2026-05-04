"""
Text content extraction utilities.
Used for AI features (question generation, AI tutor context).
For page-image rendering, see slide_renderer.py.
"""

import io
import requests


def extract_text_from_pdf(file_url_or_bytes):
    """Extract plain text from a PDF. Accepts URL string or raw bytes."""
    import fitz

    if isinstance(file_url_or_bytes, str):
        resp = requests.get(file_url_or_bytes, timeout=30)
        data = resp.content
    else:
        data = file_url_or_bytes

    doc = fitz.open(stream=data, filetype="pdf")
    pages = []
    for i, page in enumerate(doc, 1):
        pages.append({"page_number": i, "text": page.get_text()})
    doc.close()
    return pages


def extract_text_from_pptx(file_url_or_bytes):
    """Extract plain text from a PowerPoint file."""
    from pptx import Presentation

    if isinstance(file_url_or_bytes, str):
        resp = requests.get(file_url_or_bytes, timeout=30)
        data = resp.content
    else:
        data = file_url_or_bytes

    prs = Presentation(io.BytesIO(data))
    pages = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        pages.append({"page_number": i, "text": "\n".join(texts)})
    return pages


def extract_text_from_docx(file_url_or_bytes):
    """Extract plain text from a Word document."""
    from docx import Document

    if isinstance(file_url_or_bytes, str):
        resp = requests.get(file_url_or_bytes, timeout=30)
        data = resp.content
    else:
        data = file_url_or_bytes

    doc = Document(io.BytesIO(data))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [{"page_number": 1, "text": text}]


def extract_text_from_slide(file_url, file_type):
    """
    Extract text from a slide/document for AI processing.
    Returns list of {"page_number": int, "text": str}.
    """
    ft = file_type.lower()
    try:
        if 'pdf' in ft:
            return extract_text_from_pdf(file_url)
        elif 'pptx' in ft or 'presentation' in ft or 'ppt' in ft:
            return extract_text_from_pptx(file_url)
        elif 'docx' in ft or 'document' in ft or 'doc' in ft:
            return extract_text_from_docx(file_url)
        else:
            return []
    except Exception as e:
        print(f"Text extraction error ({file_type}): {e}")
        return []


def get_slide_full_text(file_url, file_type, max_chars=15000):
    """
    Return all text from a slide concatenated, capped at max_chars.
    Used as context when calling the AI.
    """
    pages = extract_text_from_slide(file_url, file_type)
    combined = "\n\n".join(p["text"] for p in pages if p["text"].strip())
    return combined[:max_chars]
